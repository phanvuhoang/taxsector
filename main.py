#!/usr/bin/env python3
"""Tax Sector Research Tool — Single-file FastAPI application"""

import os, asyncio, json, re, io, secrets, time, unicodedata
from datetime import datetime
from pathlib import Path
from typing import AsyncGenerator

import anthropic
import httpx
from fastapi import FastAPI, Request, HTTPException, Depends
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse, FileResponse, Response
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from docx import Document
from docx.shared import RGBColor
from bs4 import BeautifulSoup

# ── Config ────────────────────────────────────────────────────────────────────
PERPLEXITY_KEY = os.getenv("PERPLEXITY_API_KEY", "")
ANTHROPIC_KEY  = os.getenv("ANTHROPIC_API_KEY", "")
CLAUDE_MODEL   = os.getenv("CLAUDE_MODEL", "claude-sonnet-4-5")
APP_USER       = os.getenv("APP_USERNAME", "hoang")
APP_PASS       = os.getenv("APP_PASSWORD", "taxsector2026")
REPORTS_DIR    = Path(os.getenv("REPORTS_DIR", "./reports"))
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

# ── Auth ──────────────────────────────────────────────────────────────────────
security = HTTPBasic()

def auth(creds: HTTPBasicCredentials = Depends(security)):
    ok = (
        secrets.compare_digest(creds.username.encode(), APP_USER.encode()) and
        secrets.compare_digest(creds.password.encode(), APP_PASS.encode())
    )
    if not ok:
        raise HTTPException(401, headers={"WWW-Authenticate": "Basic"})
    return creds.username

# ── Default sections ──────────────────────────────────────────────────────────
SECTOR_SECTIONS = [
    {"id": "s1", "title": "Tong quan ve nganh",
     "sub": ["Quy mo thi truong", "Dac diem kinh doanh", "Mo hinh doanh thu/chi phi"], "enabled": True},
    {"id": "s2", "title": "Dac thu cua nganh",
     "sub": ["Chuoi cung ung upstream/downstream", "Working capital cycle", "Dac diem tai san"], "enabled": True},
    {"id": "s3", "title": "Su phat trien tai Viet Nam & Big Players",
     "sub": ["Tang truong 5 nam gan nhat", "Doanh nghiep lon nhat", "FDI", "M&A"], "enabled": True},
    {"id": "s4", "title": "Cac quy dinh phap ly quan trong",
     "sub": ["Luat chuyen nganh", "Dieu kien kinh doanh", "Giay phep", "Han che FDI"], "enabled": True},
    {"id": "s5", "title": "Phan tich cac loai thue ap dung",
     "sub": ["Thue TNDN", "Thue GTGT", "Thue Nha thau", "Thue TTDB", "Thue XNK", "Phi & le phi"], "enabled": True},
    {"id": "s6", "title": "Cac van de thue dac thu cua nganh",
     "sub": ["Rui ro doanh thu/chi phi", "Chuyen gia", "Uu dai thue",
             "Hoa don dac thu", "Khau tru thue",
             "Tranh chap thue & an le",
             "Cong van/huong dan dac thu Tong cuc Thue cho nganh"], "enabled": True},
    {"id": "s7", "title": "Thong le & van de thue quoc te",
     "sub": ["BEPS", "Chuyen gia quoc te", "So sanh voi khu vuc", "Hiep dinh thue"], "enabled": True},
]

SECTOR_SECTIONS_VI = [
    {"id": "s1", "title": "Tổng quan về ngành",
     "sub": ["Quy mô thị trường", "Đặc điểm kinh doanh", "Mô hình doanh thu/chi phí"], "enabled": True},
    {"id": "s2", "title": "Đặc thù của ngành",
     "sub": ["Chuỗi cung ứng upstream/downstream", "Working capital cycle", "Đặc điểm tài sản"], "enabled": True},
    {"id": "s3", "title": "Sự phát triển tại Việt Nam & Big Players",
     "sub": ["Tăng trưởng 5 năm gần nhất", "Doanh nghiệp lớn nhất", "FDI", "M&A"], "enabled": True},
    {"id": "s4", "title": "Các quy định pháp lý quan trọng",
     "sub": ["Luật chuyên ngành", "Điều kiện kinh doanh", "Giấy phép", "Hạn chế FDI"], "enabled": True},
    {"id": "s5", "title": "Phân tích các loại thuế áp dụng",
     "sub": ["Thuế TNDN", "Thuế GTGT", "Thuế Nhà thầu", "Thuế TTĐB", "Thuế XNK", "Phí & lệ phí"], "enabled": True},
    {"id": "s6", "title": "Các vấn đề thuế đặc thù của ngành",
     "sub": [
         "Rủi ro doanh thu/chi phí",
         "Chuyển giá",
         "Ưu đãi thuế",
         "Hóa đơn đặc thù",
         "Khấu trừ thuế",
         "Tranh chấp thuế & án lệ",
         "Công văn/hướng dẫn đặc thù của Tổng cục Thuế cho ngành",
     ], "enabled": True},
    {"id": "s7", "title": "Thông lệ & vấn đề thuế quốc tế",
     "sub": ["BEPS", "Chuyển giá quốc tế", "So sánh với khu vực", "Hiệp định thuế"], "enabled": True},
]

COMPANY_SECTIONS_VI = [
    {"id": "c1", "title": "Giới thiệu công ty",
     "sub": ["Lịch sử hình thành", "Cơ cấu sở hữu", "Hoạt động kinh doanh chính"], "enabled": True},
    {"id": "c2", "title": "Mô hình kinh doanh & chuỗi giá trị",
     "sub": ["Sản phẩm/dịch vụ", "Khách hàng mục tiêu", "Nhà cung cấp", "Chuỗi giá trị"], "enabled": True},
    {"id": "c3", "title": "Phân tích tài chính & thuế",
     "sub": ["Doanh thu & lợi nhuận", "Gánh nặng thuế", "Tỷ lệ thuế hiệu quả", "So sánh ngành"], "enabled": True},
    {"id": "c4", "title": "Rủi ro thuế đặc thù",
     "sub": [
         "Rủi ro thanh tra",
         "Chuyển giá",
         "Cấu trúc pháp lý",
         "Giao dịch liên kết",
         "Tranh chấp thuế & lịch sử thanh/kiểm tra",
         "Công văn/ruling đặc thù áp dụng cho công ty/ngành",
     ], "enabled": True},
    {"id": "c5", "title": "Khuyến nghị",
     "sub": ["Tối ưu hóa thuế", "Tuân thủ", "Cơ hội ưu đãi", "Rủi ro cần theo dõi"], "enabled": True},
]

# ── Research: Perplexity ──────────────────────────────────────────────────────
def build_query(section: dict, subject: str, mode: str) -> str:
    from datetime import datetime
    current_year = datetime.now().year
    mode_ctx = "ngành" if mode == "sector" else "công ty"

    # For company mode, add industry hint to help Perplexity find relevant laws
    subject_ctx = subject
    if mode == "company":
        subject_ctx = f"{subject} (phân tích thuế doanh nghiệp)"

    subs = ", ".join(section.get("sub", []))
    title_lower = section.get("title", "").lower()
    is_legal = any(k in title_lower for k in [
        "pháp lý","luật","quy định","thuế","thue","tài chính","rủi ro","tranh chấp"
    ])
    legal_note = (
        f"\nLƯU Ý BẮT BUỘC: Chỉ trích dẫn văn bản pháp luật CÒN HIỆU LỰC tính đến {current_year}. "
        f"Ưu tiên Luật/Nghị định/Thông tư được ban hành hoặc sửa đổi trong 2020-{current_year}. "
        f"Nếu có văn bản mới thay thế → bắt buộc dùng văn bản MỚI NHẤT, ghi rõ nó thay thế văn bản nào. "
        f"KHÔNG trích dẫn văn bản đã bị bãi bỏ."
    ) if is_legal else ""

    return (
        f"Nghiên cứu chuyên sâu về: {section['title']} — {mode_ctx} {subject_ctx} tại Việt Nam năm {current_year}\n"
        f"Chi tiết cần tìm: {subs}\n"
        f"Bao gồm: số hiệu văn bản pháp luật cụ thể, số liệu thị trường, "
        f"tên doanh nghiệp, ví dụ thực tế, nguồn đáng tin cậy"
        f"{legal_note}"
    )

async def perplexity_search(query: str, model: str = "sonar") -> dict:
    if not PERPLEXITY_KEY:
        return {"content": f"[Perplexity API key not set]\n\nQuery was: {query[:200]}", "citations": []}

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content":
             "Bạn là chuyên gia nghiên cứu thuế và pháp lý Việt Nam. "
             "Cung cấp thông tin chính xác, đầy đủ, có số liệu cụ thể và nguồn tham khảo."},
            {"role": "user", "content": query},
        ],
        "max_tokens": 2000,
        "return_citations": True,
        "return_related_questions": False,
    }

    async with httpx.AsyncClient(timeout=httpx.Timeout(90.0)) as client:
        try:
            r = await client.post(
                "https://api.perplexity.ai/chat/completions",
                headers={"Authorization": f"Bearer {PERPLEXITY_KEY}"},
                json=payload,
            )
            r.raise_for_status()
            data = r.json()
            content = data["choices"][0]["message"]["content"]
            citations = data.get("citations", [])
            return {"content": content, "citations": citations}
        except Exception as e:
            return {"content": f"[Research error: {e}]", "citations": []}

# ── Research: thuvienphapluat.vn ─────────────────────────────────────────────
TVPL_BASE = "https://thuvienphapluat.vn"

def is_legal_or_tax_section(section: dict) -> bool:
    """Return True if this section needs legal docs from thuvienphapluat.vn."""
    title = section.get("title", "").lower()
    subs  = " ".join(section.get("sub", [])).lower()
    keywords = [
        "pháp lý", "luật", "quy định", "giấy phép",
        "thuế", "thue", "tndn", "gtgt", "ttđb", "xnk",
        "nhà thầu", "chuyển giá", "ưu đãi", "tuân thủ",
        "phap ly", "van ban", "legal",
        # Company-specific additions:
        "tài chính", "gánh nặng", "tỷ lệ thuế", "rủi ro thuế",
        "thanh tra", "kiểm tra", "giao dịch liên kết", "cấu trúc pháp",
        "rui ro", "tranh chap", "ruling", "cong van",
    ]
    return any(kw in title + " " + subs for kw in keywords)

def build_tvpl_query(section: dict, subject: str, mode: str = "sector") -> str:
    title = section.get("title", "").lower()
    subs  = " ".join(section.get("sub", [])).lower()

    # For company mode: add "doanh nghiệp" context
    base_subject = subject
    if mode == "company":
        base_subject = f"{subject} doanh nghiệp"

    kw_map = {
        "thuế gtgt": "thuế giá trị gia tăng",
        "thuế tndn": "thuế thu nhập doanh nghiệp",
        "thuế ttđb": "thuế tiêu thụ đặc biệt",
        "thuế xnk": "thuế xuất nhập khẩu",
        "nhà thầu": "thuế nhà thầu",
        "chuyển giá": "chuyển giá giao dịch liên kết",
        "ưu đãi": f"ưu đãi thuế {base_subject}",
        "tranh chấp": f"tranh chấp thuế {base_subject}",
        "giao dịch liên kết": "chuyển giá giao dịch liên kết",
        "tài chính": f"thuế {base_subject}",
        "thuế": f"thuế {base_subject}",
        "pháp lý": base_subject,
        "luật": base_subject,
        "quy định": base_subject,
    }
    for kw, q in kw_map.items():
        if kw in title or kw in subs:
            return q
    return f"thuế {base_subject}"

async def tvpl_search(query: str, max_results: int = 10) -> list:
    """Scrape thuvienphapluat.vn for currently-in-effect legal documents."""
    params = {
        "q": query,
        "sbt": "1",   # sort by date descending (newest first)
        "efts": "1",  # còn hiệu lực only — filters out superseded docs
        "page": "1",
    }
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        ),
        "Accept-Language": "vi-VN,vi;q=0.9,en;q=0.8",
        "Referer": TVPL_BASE,
    }

    async with httpx.AsyncClient(timeout=httpx.Timeout(30.0), follow_redirects=True) as client:
        try:
            r = await client.get(
                f"{TVPL_BASE}/van-ban-phap-luat.aspx",
                params=params,
                headers=headers,
            )
            r.raise_for_status()
        except Exception as e:
            return [{"error": str(e)}]

    soup = BeautifulSoup(r.text, "lxml")
    results = []

    items = (
        soup.select("div.doc-item") or
        soup.select("ul.result-list > li") or
        soup.select(".document-list .item") or
        soup.select("table.list-vb tr") or
        soup.select("a[href*='/van-ban/']")
    )

    doc_types = [
        "Luật", "Nghị định", "Thông tư", "Quyết định",
        "Nghị quyết", "Chỉ thị", "Công văn", "Pháp lệnh",
        "Hiệp định", "Thông tư liên tịch",
    ]
    issuers = [
        "Quốc hội", "Chính phủ", "Bộ Tài chính", "Bộ Kế hoạch và Đầu tư",
        "Tổng cục Thuế", "Bộ Công Thương", "Ngân hàng Nhà nước",
        "Bộ Lao động", "UBND", "Bộ Xây dựng",
    ]

    for item in items[:max_results]:
        doc = {}
        link = item.select_one("a[href*='/van-ban/']") or (item if item.name == "a" else None)
        if not link:
            continue
        doc["title"] = link.get_text(strip=True)
        href = link.get("href", "")
        doc["url"] = href if href.startswith("http") else f"{TVPL_BASE}{href}"
        if not doc["title"] or not doc["url"]:
            continue
        meta = item.get_text(" ", strip=True)
        doc["doc_type"] = next((t for t in doc_types if t in doc["title"]), "Văn bản")
        date_m = re.search(r'\b(\d{2}/\d{2}/\d{4})\b', meta)
        doc["issued_date"] = date_m.group(1) if date_m else ""
        doc["issuer"] = next((i for i in issuers if i in meta), "")
        doc["status"] = "Còn hiệu lực"
        results.append(doc)

    return results

def format_tvpl_results(docs: list) -> str:
    """Format TVPL docs as context text for Claude."""
    valid = [d for d in docs if "error" not in d and d.get("title")]
    if not valid:
        return ""
    lines = ["=== VĂN BẢN PHÁP LUẬT HIỆN HÀNH (nguồn: thuvienphapluat.vn) ===\n"]
    for i, d in enumerate(valid, 1):
        lines.append(
            f"{i}. [{d['doc_type']}] {d['title']}\n"
            f"   URL: {d['url']}\n"
            f"   Ban hành: {d.get('issued_date','')} | "
            f"Cơ quan: {d.get('issuer','')} | Trạng thái: Còn hiệu lực\n"
        )
    return "\n".join(lines)

async def _empty():
    return []

# ── Context filtering ─────────────────────────────────────────────────────────
KEYWORD_MAP = {
    "pháp lý": ["s4", "c4"],
    "luật": ["s4", "c4"],
    "quy định": ["s4", "c4"],
    "thuế": ["s5", "s6", "s7", "c3", "c4"],
    "tổng quan": ["s1", "c1"],
    "đặc thù": ["s2", "s6"],
    "big player": ["s3"],
    "doanh nghiep": ["s3", "c1"],
    "quoc te": ["s7"],
    "tài chính": ["c3"],
    "khuyến nghị": ["c5"],
    "mô hình": ["s1", "s2", "c2"],
    "chuỗi": ["s2", "c2"],
}

def filter_context(all_results: dict, section: dict) -> str:
    title_lower = section["title"].lower()
    relevant = {section["id"]}
    for kw, ids in KEYWORD_MAP.items():
        if kw in title_lower:
            relevant.update(ids)

    parts = []
    # Own section first
    if section["id"] in all_results:
        parts.append(all_results[section["id"]].get("content", "")[:5000])
    # Related sections (shorter)
    for sid in relevant:
        if sid != section["id"] and sid in all_results:
            parts.append(all_results[sid].get("content", "")[:2000])

    return "\n\n---\n\n".join(parts)[:10000]

# ── Claude per-section streaming ──────────────────────────────────────────────
def build_section_prompt(section: dict, subject: str, context: str, mode: str, num: int) -> str:
    mode_ctx = "ngành" if mode == "sector" else "công ty"
    sub_list = "\n".join(f"- {s}" for s in section.get("sub", []))
    title = section["title"]

    is_legal = any(k in title.lower() for k in ["pháp lý", "luật", "quy định", "phap ly", "luat"])
    is_tax   = any(k in title.lower() for k in ["thuế", "thue"])

    table_block = ""
    if is_legal or is_tax:
        table_block = """
QUAN TRỌNG — VĂN BẢN PHÁP LUẬT:
1. CHỈ sử dụng văn bản đang CÒN HIỆU LỰC tại thời điểm hiện tại (2025-2026).
2. Nếu có văn bản mới thay thế/sửa đổi văn bản cũ → chỉ trích dẫn văn bản MỚI NHẤT.
3. KHÔNG trích dẫn văn bản đã bị bãi bỏ, thay thế hoặc hết hiệu lực.
4. Ưu tiên văn bản từ phần "VĂN BẢN PHÁP LUẬT HIỆN HÀNH" trong dữ liệu nghiên cứu.
5. Bắt đầu section bằng bảng tổng hợp văn bản (ít nhất 6-8 văn bản còn hiệu lực):
<table>
  <thead><tr>
    <th>Số hiệu</th><th>Tên văn bản</th><th>Loại</th>
    <th>Ngày ban hành</th><th>Hiệu lực</th><th>Ghi chú sửa đổi</th>
  </tr></thead>
  <tbody>... điền thực tế, ghi rõ nếu văn bản này thay thế văn bản nào ...</tbody>
</table>
"""

    return f"""Bạn là chuyên gia thuế Big 4 (Deloitte/PwC/EY/KPMG) viết báo cáo phân tích thuế chuyên nghiệp bằng tiếng Việt.

Viết PHẦN {num}: "{title}" trong báo cáo phân tích về {mode_ctx}: **{subject}**

Các chủ đề bắt buộc đề cập:
{sub_list}
{table_block}
DỮ LIỆU NGHIÊN CỨU (dùng thông tin này, bổ sung thêm kiến thức của bạn):
{context}

YÊU CẦU TUYỆT ĐỐI:
1. Output là HTML THUẦN TÚY — KHÔNG markdown, KHÔNG backticks, KHÔNG code block
2. Bắt đầu bằng: <h2>{num}. {title}</h2>
3. Dùng thẻ HTML: <h3>, <p>, <ul><li>, <ol><li>, <table> — KHÔNG <div> thừa
4. Văn phong: chuyên nghiệp, cụ thể, dẫn chứng số liệu và tên văn bản pháp luật thực tế
5. TRÍCH DẪN NGUỒN — BẮT BUỘC TUYỆT ĐỐI:
   - Sau MỖI câu có số liệu, tên văn bản, hoặc thông tin cụ thể → chèn ngay: <a href="URL" target="_blank" rel="noopener">[N]</a>
   - N là số thứ tự tăng dần từ 1 trong toàn bộ phần này
   - URL phải là URL thực tế từ dữ liệu nghiên cứu (Perplexity citations hoặc thuvienphapluat.vn)
   - Nếu không có URL cụ thể cho câu đó → dùng URL tổng quát của nguồn
   - KHÔNG viết [N] mà không có thẻ <a href=...>
   - KHÔNG gộp nhiều câu dùng chung 1 citation số
   Ví dụ đúng: Thuế GTGT hiện hành là 10%.<a href="https://thuvienphapluat.vn/van-ban/..." target="_blank" rel="noopener">[1]</a>
   Ví dụ sai: Thuế GTGT hiện hành là 10%.[1] hoặc [1] không có href
6. Tối thiểu 700 từ — đầy đủ, không rút gọn
7. KHÔNG viết lời mở đầu/kết luận tổng quát — chỉ nội dung của phần này"""

async def claude_stream_section(
    section: dict, subject: str, context: str, mode: str, num: int
) -> AsyncGenerator[str, None]:
    if not ANTHROPIC_KEY:
        yield f"<h2>{num}. {section['title']}</h2><p><em>[Anthropic API key not configured]</em></p>"
        return

    prompt = build_section_prompt(section, subject, context, mode, num)
    client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_KEY)

    try:
        async with client.messages.stream(
            model=CLAUDE_MODEL,
            max_tokens=8192,
            messages=[{"role": "user", "content": prompt}],
        ) as stream:
            async for text in stream.text_stream:
                yield text
    except Exception as e:
        yield f'<p style="color:red">[Error in section {num}: {e}]</p>'

# ── Report persistence ────────────────────────────────────────────────────────
CAVEAT_HTML = """
<div style="margin-top:3rem;padding:1rem 1.5rem;background:#f8fafc;border-top:2px solid #e2e8f0;
            border-radius:.5rem;font-size:.8rem;color:#64748b;line-height:1.6">
  <strong>&#9888;&#65039; Lưu ý quan trọng:</strong> Báo cáo này được tạo tự động bởi
  <strong>Tax Sector Research AI</strong> dựa trên dữ liệu từ Perplexity (sonar model)
  và thuvienphapluat.vn. Nội dung mang tính tham khảo, không thay thế tư vấn pháp lý
  hoặc thuế chuyên nghiệp. Người dùng cần kiểm chứng độc lập trước khi áp dụng.
  Thông tin pháp luật có thể thay đổi — vui lòng xác nhận hiệu lực văn bản tại
  <a href="https://thuvienphapluat.vn" target="_blank" rel="noopener">thuvienphapluat.vn</a>.
  <br><em>Ngày tạo: {date_str}</em>
</div>
"""

def safe_filename(s: str) -> str:
    return re.sub(r'[\\/*?:"<>|]', "", s)[:60].strip()

def save_report(subject: str, html_content: str) -> str:
    now = datetime.now()
    date_str = now.strftime("%d%m%Y")
    time_str = now.strftime("%H%M")
    base = safe_filename(subject)
    name = f"{date_str} - {base} - {time_str}.html"

    # Handle rare collision (same subject, same minute)
    if (REPORTS_DIR / name).exists():
        name = f"{date_str} - {base} - {now.strftime('%H%M%S')}.html"

    full_html = f"""<!DOCTYPE html>
<html lang="vi">
<head>
<meta charset="UTF-8">
<title>Phan Tich Thue — {subject}</title>
<style>
body{{font-family:system-ui,sans-serif;max-width:900px;margin:2rem auto;padding:0 1.5rem;line-height:1.75;color:#1e293b}}
h1{{color:#028a39;border-bottom:3px solid #028a39;padding-bottom:.5rem}}
h2{{color:#028a39;border-bottom:2px solid #028a39;padding-bottom:4px;margin-top:2.5rem}}
h3{{margin-top:1.5rem;color:#1e293b}}
table{{width:100%;border-collapse:collapse;margin:1rem 0;font-size:.875rem}}
th{{background:#028a39;color:#fff;padding:8px;text-align:left}}
td{{padding:6px 8px;border:1px solid #e2e8f0}}
tr:nth-child(even) td{{background:#f8fafc}}
a{{color:#028a39}}
ul,ol{{padding-left:1.5rem}}
li{{margin:.25rem 0}}
p{{margin:.6rem 0}}
@media print{{body{{max-width:100%}}}}
</style>
</head>
<body>
<h1>Phan Tich Thue — {subject}</h1>
<p><em>Ngay tao: {now.strftime("%d/%m/%Y %H:%M")}</em></p>
<div id="report-body">
{html_content}
</div>
{CAVEAT_HTML.format(date_str=now.strftime("%d/%m/%Y %H:%M"))}
</body>
</html>"""
    (REPORTS_DIR / name).write_text(full_html, encoding="utf-8")
    return name

def linkify_citations(html: str, citations: list) -> str:
    """Safety net: replace bare [N] with linked version using citation URLs."""
    def replacer(m):
        idx = int(m.group(1)) - 1
        if 0 <= idx < len(citations) and citations[idx]:
            url = citations[idx]
            return f'<a href="{url}" target="_blank" rel="noopener">[{idx+1}]</a>'
        return m.group(0)
    # Only replace [N] that are NOT already inside an <a> tag
    return re.sub(r'(?<!href=")\[(\d+)\](?![^<]*</a>)', replacer, html)

# ── FastAPI app ───────────────────────────────────────────────────────────────
app = FastAPI(title="Tax Sector Research")

# ── SSE stream endpoint ───────────────────────────────────────────────────────
@app.post("/stream")
async def stream_report(request: Request, _user: str = Depends(auth)):
    body = await request.json()
    subject  = body.get("subject", "").strip()
    mode     = body.get("mode", "sector")
    sections = body.get("sections", [])
    sonar    = body.get("sonar_model", "sonar")

    if not subject:
        raise HTTPException(400, "Missing subject")

    enabled = [s for s in sections if s.get("enabled")]
    if not enabled:
        raise HTTPException(400, "No sections enabled")

    async def generate():
        def sse(data: dict) -> str:
            return f"data: {json.dumps(data, ensure_ascii=False)}\n\n"

        yield sse({"type": "ping"})
        yield sse({"type": "status", "message": f"Đang nghiên cứu '{subject}'..."})

        # Phase 1: parallel research — Perplexity + thuvienphapluat.vn
        all_results: dict = {}
        all_tvpl: dict = {}
        all_citations: list = []
        total = len(enabled)

        for batch_start in range(0, total, 4):
            batch = enabled[batch_start:batch_start + 4]

            for i, sec in enumerate(batch):
                yield sse({
                    "type": "progress",
                    "step": batch_start + i + 1,
                    "total": total,
                    "label": f"Đang nghiên cứu: {sec['title']}",
                })

            perplexity_tasks = [
                perplexity_search(build_query(s, subject, mode), sonar)
                for s in batch
            ]
            tvpl_tasks = [
                tvpl_search(build_tvpl_query(s, subject, mode))
                if is_legal_or_tax_section(s)
                else _empty()
                for s in batch
            ]

            perplexity_results, tvpl_results = await asyncio.gather(
                asyncio.gather(*perplexity_tasks),
                asyncio.gather(*tvpl_tasks),
            )

            for sec, pres, tres in zip(batch, perplexity_results, tvpl_results):
                all_results[sec["id"]] = pres
                all_tvpl[sec["id"]]    = tres or []
                all_citations.extend(pres.get("citations", []))
                all_citations.extend(
                    d.get("url", "") for d in (tres or []) if d.get("url")
                )
                tvpl_note = f" (+{len(tres)} văn bản PL)" if tres else ""
                yield sse({
                    "type": "progress",
                    "step": batch_start + batch.index(sec) + 1,
                    "total": total,
                    "label": f"Xong research: {sec['title']}{tvpl_note}",
                })

        # Phase 2: Claude per section
        yield sse({"type": "ai_start", "total": len(enabled)})
        full_html = ""

        for i, section in enumerate(enabled):
            yield sse({
                "type": "status",
                "message": f"AI đang viết phần {i + 1}/{len(enabled)}: {section['title']}",
            })
            ctx = filter_context(all_results, section)

            # Prepend TVPL legal docs for law/tax sections
            tvpl_docs = all_tvpl.get(section["id"], [])
            if tvpl_docs and is_legal_or_tax_section(section):
                tvpl_text = format_tvpl_results(tvpl_docs)
                if tvpl_text:
                    ctx = tvpl_text + "\n\n" + ctx

            # Inject citation URLs so Claude can use real links
            section_citations = all_results.get(section["id"], {}).get("citations", [])
            tvpl_urls = [d.get("url","") for d in tvpl_docs if d.get("url")]
            all_section_urls = section_citations + tvpl_urls
            if all_section_urls:
                url_list = "\n".join(f"[{i+1}] {u}" for i,u in enumerate(all_section_urls[:20]))
                ctx = f"=== DANH SÁCH URL NGUỒN (dùng cho citation) ===\n{url_list}\n\n" + ctx

            sec_html = ""

            async for chunk in claude_stream_section(section, subject, ctx, mode, i + 1):
                yield sse({"type": "chunk", "text": chunk})
                sec_html += chunk
                await asyncio.sleep(0)

            full_html += sec_html + "\n"
            yield sse({"type": "ping"})

        # Citations
        unique_urls = list(dict.fromkeys(all_citations))
        yield sse({"type": "citations", "urls": unique_urls})

        # Linkify any bare [N] citations
        unique_citations = list(dict.fromkeys(all_citations))
        full_html = linkify_citations(full_html, unique_citations)

        # Save — append citations to report before saving
        try:
            if unique_urls:
                refs_html = '<hr><h2>Nguồn tham khảo</h2><ol>'
                for url in unique_urls[:50]:
                    refs_html += f'<li><a href="{url}" target="_blank" rel="noopener">{url}</a></li>'
                refs_html += '</ol>'
                full_html_with_refs = full_html + refs_html
            else:
                full_html_with_refs = full_html
            filename = save_report(subject, full_html_with_refs)
        except Exception:
            filename = None

        yield sse({"type": "done", "filename": filename, "drive": False})

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )

# ── Default sections ──────────────────────────────────────────────────────────
@app.get("/default-sections")
def default_sections(mode: str = "sector", _user: str = Depends(auth)):
    return SECTOR_SECTIONS_VI if mode == "sector" else COMPANY_SECTIONS_VI

# ── Suggest subsections ───────────────────────────────────────────────────────
@app.post("/suggest-subsections")
async def suggest_subsections(request: Request, _user: str = Depends(auth)):
    body = await request.json()
    title   = body.get("title", "")
    subject = body.get("subject", "")

    if not ANTHROPIC_KEY:
        return {"suggestions": []}

    prompt = (
        f'Đề xuất 4-5 chủ đề con phù hợp cho phần "{title}" '
        f'trong báo cáo phân tích thuế về: {subject}\n'
        f'Trả về JSON array, ví dụ: ["Chủ đề 1", "Chủ đề 2"]\n'
        f'Chỉ trả về JSON array, không giải thích thêm.'
    )
    client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_KEY)
    try:
        msg = await client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=300,
            messages=[{"role": "user", "content": prompt}],
        )
        content = msg.content[0].text
        match = re.search(r'\[.*\]', content, re.DOTALL)
        suggestions = json.loads(match.group()) if match else []
    except Exception:
        suggestions = []
    return {"suggestions": suggestions}

# ── Reports ───────────────────────────────────────────────────────────────────
@app.get("/reports")
def list_reports(_user: str = Depends(auth)):
    reports = []
    for f in sorted(REPORTS_DIR.glob("*.html"), key=lambda x: x.stat().st_mtime, reverse=True):
        stat = f.stat()
        size = f"{stat.st_size // 1024} KB" if stat.st_size >= 1024 else f"{stat.st_size} B"
        reports.append({
            "name": f.name,
            "date": datetime.fromtimestamp(stat.st_mtime).strftime("%d/%m/%Y %H:%M"),
            "size": size,
        })
    return reports

@app.get("/report/{name:path}")
def get_report(name: str, _user: str = Depends(auth)):
    path = (REPORTS_DIR / name).resolve()
    try:
        path.relative_to(REPORTS_DIR.resolve())
    except ValueError:
        raise HTTPException(400, "Invalid path")
    if not path.exists():
        raise HTTPException(404)
    return HTMLResponse(path.read_text(encoding="utf-8"))

@app.delete("/report/{name:path}")
def delete_report(name: str, _user: str = Depends(auth)):
    path = (REPORTS_DIR / name).resolve()
    try:
        path.relative_to(REPORTS_DIR.resolve())
    except ValueError:
        raise HTTPException(400, "Invalid path")
    if not path.exists():
        raise HTTPException(404)
    path.unlink()
    return {"ok": True}

# ── DOCX export ───────────────────────────────────────────────────────────────
def normalize_text(text: str) -> str:
    """NFC normalize để fix lỗi latin-1 encoding trong python-docx."""
    return unicodedata.normalize('NFC', text)

@app.post("/docx")
async def export_docx(request: Request, _user: str = Depends(auth)):
    try:
        body    = await request.json()
        html    = body.get("html", "")
        subject = body.get("subject", "Báo cáo")

        if not html.strip():
            raise HTTPException(400, "Nội dung báo cáo trống")

        doc = Document()
        # Fix encoding: set core properties with ASCII-safe values
        doc.core_properties.author = "Tax Sector Research"
        doc.core_properties.title = ""
        doc.core_properties.subject = ""
        doc.core_properties.keywords = ""
        doc.core_properties.description = ""
        title_para = doc.add_heading(normalize_text(f"Phân Tích Thuế — {subject}"), 0)
        if title_para.runs:
            title_para.runs[0].font.color.rgb = RGBColor(0x02, 0x8A, 0x39)
        doc.add_paragraph(normalize_text(f"Ngày tạo: {datetime.now().strftime('%d/%m/%Y')}"))

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup.find_all(["a", "strong", "em", "span", "b", "i"]):
            tag.insert_before(" ")
            tag.insert_after(" ")

        for el in soup.find_all(["h2", "h3", "p", "li", "table"]):
            try:
                text = normalize_text(" ".join(el.get_text(" ", strip=False).split()))
                if not text:
                    continue
                tag = el.name
                if tag == "h2":
                    p = doc.add_heading(text, level=1)
                    if p.runs:
                        p.runs[0].font.color.rgb = RGBColor(0x02, 0x8A, 0x39)
                elif tag == "h3":
                    doc.add_heading(text, level=2)
                elif tag == "p":
                    doc.add_paragraph(text)
                elif tag == "li":
                    try:
                        doc.add_paragraph(text, style="List Bullet")
                    except KeyError:
                        doc.add_paragraph(f"• {text}")
                elif tag == "table":
                    rows = el.find_all("tr")
                    if not rows:
                        continue
                    cols = max(len(r.find_all(["th", "td"])) for r in rows)
                    if cols == 0:
                        continue
                    t = doc.add_table(rows=0, cols=cols)
                    try:
                        t.style = "Table Grid"
                    except KeyError:
                        pass
                    for row in rows:
                        cells = row.find_all(["th", "td"])
                        row_cells = t.add_row().cells
                        for j, cell in enumerate(cells[:cols]):
                            row_cells[j].text = normalize_text(cell.get_text(strip=True))
            except Exception as el_err:
                print(f"[DOCX] Skip element {el.name}: {el_err}")
                continue

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        safe_ascii = re.sub(r'[^\x01-\x7F]', '', subject)[:50].strip().replace(' ', '_') or 'BaoCao'
        from urllib.parse import quote as _quote
        safe_utf8 = _quote(f'PhanTichThue_{subject[:50]}.docx', safe='')
        cd = f'attachment; filename="PhanTichThue_{safe_ascii}.docx"; filename*=UTF-8''{safe_utf8}'
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": cd},
        )

    except HTTPException:
        raise
    except Exception as e:
        print(f"[DOCX] Export failed: {e}")
        raise HTTPException(500, f"Xuất DOCX thất bại: {str(e)[:200]}")

# ── PPTX export ───────────────────────────────────────────────────────────────
@app.post("/slides")
async def export_pptx(request: Request, _user: str = Depends(auth)):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    body    = await request.json()
    html    = body.get("html", "")
    subject = body.get("subject", "Báo cáo")

    BRAND = PptxRGB(0x02, 0x8A, 0x39)
    WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
    DARK  = PptxRGB(0x1E, 0x29, 0x3B)
    LIGHT = PptxRGB(0xF8, 0xFA, 0xFC)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_rect(slide, l, t, w, h, fill_color=None):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        shape.line.fill.background()
        return shape

    def add_text_box(slide, text, l, t, w, h, font_size=18, bold=False,
                     color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    soup = BeautifulSoup(html, "html.parser")

    # Slide 1: Cover
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=BRAND)
    add_text_box(slide, "PHÂN TÍCH THUẾ",
                 0.8, 1.5, 11.5, 1.0,
                 font_size=20, color=PptxRGB(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)
    add_text_box(slide, subject,
                 0.8, 2.5, 11.5, 2.0,
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, datetime.now().strftime("%d/%m/%Y"),
                 0.8, 5.5, 11.5, 0.8,
                 font_size=16, color=PptxRGB(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)

    # Content slides: one per h2
    for h2 in soup.find_all("h2"):
        title_text = h2.get_text(strip=True)
        bullets = []
        for sib in h2.find_next_siblings():
            if sib.name == "h2":
                break
            if sib.name in ("p", "li"):
                t = sib.get_text(strip=True)
                if t and len(t) > 10:
                    bullets.append(t[:200])
            elif sib.name in ("ul", "ol"):
                for li in sib.find_all("li"):
                    t = li.get_text(strip=True)
                    if t:
                        bullets.append(t[:200])
            if len(bullets) >= 7:
                break

        slide = prs.slides.add_slide(blank_layout)
        add_rect(slide, 0, 0, 13.333, 1.4, fill_color=BRAND)
        add_text_box(slide, title_text,
                     0.4, 0.1, 12.5, 1.2,
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_rect(slide, 0, 1.4, 13.333, 6.1, fill_color=LIGHT)

        if bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5)
            )
            txBox.word_wrap = True
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets[:7]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(6)
                run = p.add_run()
                run.text = f"▪  {bullet}"
                run.font.size = Pt(16)
                run.font.color.rgb = DARK
        else:
            add_text_box(slide, "(Xem báo cáo đầy đủ để biết chi tiết)",
                         0.5, 2.5, 12.3, 1.0,
                         font_size=16, color=PptxRGB(0x94, 0xA3, 0xB8))

    # Last slide: Thank you
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=BRAND)
    add_text_box(slide, "Cảm ơn",
                 0.8, 2.5, 11.5, 1.5,
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Báo cáo được tạo bởi Tax Sector Research AI",
                 0.8, 4.2, 11.5, 1.0,
                 font_size=16, color=PptxRGB(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe = re.sub(r'[^\w\s-]', '', subject)[:50].strip().replace(' ', '_')
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="TaxSlides_{safe}.pptx"'},
    )

# ── Health ────────────────────────────────────────────────────────────────────
@app.get("/health")
def health(_user: str = Depends(auth)):
    return {
        "status": "ok",
        "model": CLAUDE_MODEL,
        "anthropic_configured": bool(ANTHROPIC_KEY),
        "perplexity_configured": bool(PERPLEXITY_KEY),
        "tvpl_scraper": "enabled",
    }

# ── Frontend ──────────────────────────────────────────────────────────────────
HTML_PAGE = r"""<!DOCTYPE html>
<html lang="vi" class="scroll-smooth">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Tax Sector Research</title>
<link rel="icon" type="image/svg+xml" href="/favicon.svg">
<link rel="shortcut icon" href="/favicon.svg">
<script src="https://cdn.tailwindcss.com"></script>
<style>
:root{--brand:#028a39;--brand-dk:#016d2d;--bg:#f8fafc;--surface:#fff;--border:#e2e8f0;--text:#1e293b;--muted:#64748b}
body.dark{--bg:#0f172a;--surface:#1e293b;--border:#334155;--text:#e2e8f0;--muted:#94a3b8}
body{background:var(--bg);color:var(--text);font-family:system-ui,sans-serif;transition:background .2s,color .2s}

/* Reading progress */
#reading-bar{position:fixed;top:0;left:0;height:3px;background:var(--brand);width:0;z-index:1000;transition:width .1s}

/* Section card */
.sec-card{background:var(--surface);border:1px solid var(--border);border-radius:.5rem;padding:.75rem;margin-bottom:.5rem;transition:opacity .2s}
.sec-card.off{opacity:.45}
.sub-chip{display:inline-flex;align-items:center;gap:.2rem;background:#f0fdf4;border:1px solid #bbf7d0;border-radius:9999px;padding:.1rem .5rem;font-size:.75rem}
body.dark .sub-chip{background:#064e3b;border-color:#065f46}
.sub-chip button{color:#9ca3af;cursor:pointer;line-height:1;font-size:.85rem}
.sub-chip button:hover{color:#ef4444}

/* Report content */
#report-content h2{font-size:1.2rem;font-weight:700;color:var(--brand);margin:2rem 0 .75rem;border-bottom:2px solid var(--brand);padding-bottom:.25rem}
#report-content h3{font-size:1rem;font-weight:600;margin:1.25rem 0 .4rem}
#report-content p{margin:.5rem 0;line-height:1.75}
#report-content ul,#report-content ol{padding-left:1.5rem;margin:.5rem 0}
#report-content li{margin:.25rem 0}
#report-content table{width:100%;border-collapse:collapse;margin:1rem 0;font-size:.875rem;display:block;overflow-x:auto}
#report-content th{background:var(--brand);color:#fff;padding:.5rem .6rem;text-align:left}
#report-content td{padding:.4rem .6rem;border:1px solid var(--border)}
#report-content tr:nth-child(even) td{background:#f8fafc}
body.dark #report-content tr:nth-child(even) td{background:#162032}
#report-content a{color:var(--brand)}

/* Surface/border helpers */
.surface{background:var(--surface);border:1px solid var(--border)}
.btn{display:inline-flex;align-items:center;gap:.375rem;padding:.375rem .75rem;border-radius:.5rem;font-size:.875rem;font-weight:500;cursor:pointer;transition:background .15s}
.btn-gray{background:#f1f5f9;color:var(--text)}
.btn-gray:hover{background:#e2e8f0}
.btn-green{background:#f0fdf4;border:1px solid #bbf7d0;color:#15803d}
.btn-green:hover{background:#dcfce7}
body.dark .btn-gray{background:#334155;color:var(--text)}
body.dark .btn-gray:hover{background:#475569}
body.dark .btn-green{background:#064e3b;border-color:#065f46;color:#6ee7b7}

/* Fade in */
.fade-in{animation:fi .25s ease}
@keyframes fi{from{opacity:0;transform:translateY(6px)}to{opacity:1;transform:none}}

/* Print */
@media print{.no-print{display:none!important}#report-content{font-size:11pt}}

/* Input dark */
body.dark input,body.dark textarea,body.dark select{background:#1e293b;color:var(--text);border-color:var(--border)}
body.dark .sec-card input[type=text]{background:transparent}
</style>
</head>
<body>
<div id="reading-bar"></div>

<!-- ── Login Modal ──────────────────────────────────────────── -->
<div id="login-modal" class="fixed inset-0 bg-black/60 flex items-center justify-center z-50">
  <div class="bg-white rounded-2xl shadow-2xl p-8 w-full max-w-sm mx-4">
    <div class="text-center mb-6">
      <div class="text-5xl mb-3">📊</div>
      <h1 class="text-2xl font-bold text-gray-800">Tax Research</h1>
      <p class="text-gray-400 text-sm mt-1">Phân tích thuế AI — Việt Nam</p>
    </div>
    <div class="space-y-3">
      <input id="li-user" type="text" placeholder="Tên đăng nhập" value="hoang"
        class="w-full border rounded-lg px-4 py-2.5 focus:outline-none focus:ring-2 focus:ring-green-400 text-gray-800">
      <input id="li-pass" type="password" placeholder="Mật khẩu"
        class="w-full border rounded-lg px-4 py-2.5 focus:outline-none focus:ring-2 focus:ring-green-400 text-gray-800"
        onkeydown="if(event.key==='Enter')doLogin()">
      <button onclick="doLogin()"
        class="w-full bg-green-600 hover:bg-green-700 text-white font-semibold py-2.5 rounded-lg transition">
        Đăng nhập
      </button>
      <p id="li-err" class="text-red-500 text-sm text-center hidden">Sai tên đăng nhập hoặc mật khẩu</p>
    </div>
  </div>
</div>

<!-- ── App Shell ─────────────────────────────────────────────── -->
<div id="app" class="max-w-4xl mx-auto px-4 py-5">

  <!-- Header -->
  <header class="flex items-center justify-between mb-5">
    <div class="flex items-center gap-3 cursor-pointer" onclick="newReport()">
      <span class="text-3xl">📊</span>
      <div>
        <h1 class="text-lg font-bold">Tax Sector Research</h1>
        <p class="text-xs" style="color:var(--muted)">Phân tích thuế AI — Việt Nam</p>
      </div>
    </div>
    <button onclick="toggleDark()" id="dark-btn" class="text-xl p-2 rounded-lg btn-gray">☀️</button>
  </header>

  <!-- ── SETUP PHASE ─────────────────────────────────────────── -->
  <div id="ph-setup">
    <!-- Tabs -->
    <div class="flex gap-2 mb-4">
      <button id="tab-sector" onclick="setMode('sector')"
        class="px-5 py-2 rounded-full text-sm font-medium transition bg-green-600 text-white">
        🏭 Ngành / Lĩnh vực
      </button>
      <button id="tab-company" onclick="setMode('company')"
        class="px-5 py-2 rounded-full text-sm font-medium transition bg-gray-100 text-gray-600 hover:bg-gray-200">
        🏢 Doanh nghiệp
      </button>
    </div>

    <!-- Reports shortcut -->
    <div class="mt-4 text-center">
      <button onclick="openReports()"
        class="btn btn-gray text-sm px-5 py-2">
        📂 Báo cáo đã lưu
        <span id="reports-badge" class="ml-1 text-xs opacity-60"></span>
      </button>
    </div>

    <!-- Subject input -->
    <div class="surface rounded-xl p-5 mb-4 shadow-sm">
      <label id="subj-label" class="block text-sm font-medium mb-2">Tên ngành / lĩnh vực:</label>
      <input id="subj-input" type="text"
        placeholder="VD: Bất động sản, Fintech, Sản xuất thép, Thương mại điện tử..."
        class="w-full border border-gray-200 rounded-lg px-4 py-3 text-lg focus:outline-none focus:ring-2 focus:ring-green-400"
        onkeydown="if(event.key==='Enter')startResearch()">
      <div class="flex items-center gap-5 mt-3">
        <span class="text-sm" style="color:var(--muted)">Research model:</span>
        <label class="flex items-center gap-1.5 cursor-pointer text-sm">
          <input type="radio" name="sonar" value="sonar" checked class="accent-green-600"> Sonar (nhanh)
        </label>
        <label class="flex items-center gap-1.5 cursor-pointer text-sm">
          <input type="radio" name="sonar" value="sonar-pro" class="accent-green-600"> Sonar Pro (sâu)
        </label>
      </div>
    </div>

    <!-- Section editor -->
    <div class="surface rounded-xl p-5 mb-4 shadow-sm">
      <div class="flex items-center justify-between mb-3">
        <h2 class="font-semibold text-sm uppercase tracking-wide" style="color:var(--muted)">Cấu trúc báo cáo</h2>
        <div class="flex gap-2">
          <button onclick="resetSections()" class="btn btn-gray text-xs">↩ Đặt lại</button>
          <button onclick="addSection()" class="btn btn-green text-xs">+ Thêm phần</button>
        </div>
      </div>
      <div id="sections-list"></div>
    </div>

    <!-- Submit -->
    <button onclick="startResearch()"
      class="w-full bg-green-600 hover:bg-green-700 active:bg-green-800 text-white font-bold
             py-4 rounded-xl text-lg shadow-md transition">
      🔍 Bắt đầu phân tích
    </button>
    <p class="text-center text-xs mt-2" style="color:var(--muted)">Ước tính 2–5 phút tuỳ số lượng phần</p>
  </div>

  <!-- ── PROGRESS PHASE ──────────────────────────────────────── -->
  <div id="ph-progress" class="hidden">
    <div class="surface rounded-xl p-5 mb-4 shadow-sm">
      <div class="flex items-center gap-3 mb-4">
        <div class="animate-spin text-2xl select-none">⚙️</div>
        <div class="flex-1">
          <p id="prog-status" class="font-medium">Đang khởi động...</p>
          <p id="prog-subject" class="text-sm" style="color:var(--muted)"></p>
        </div>
        <button onclick="cancelResearch()" class="btn btn-gray text-xs text-red-500">✕ Huỷ</button>
      </div>
      <div class="h-2 bg-gray-100 rounded-full overflow-hidden">
        <div id="prog-bar" class="h-full bg-green-500 transition-all duration-500 rounded-full" style="width:0%"></div>
      </div>
      <div class="flex justify-between text-xs mt-1" style="color:var(--muted)">
        <span id="prog-label">Bắt đầu...</span>
        <span id="prog-pct">0%</span>
      </div>
    </div>
    <div id="step-list" class="space-y-2 mb-4"></div>
    <div id="partial-wrap" class="surface rounded-xl p-5 hidden">
      <p class="font-medium text-green-700 mb-3 text-sm">📝 Đang tạo báo cáo...</p>
      <div id="partial-body" class="text-sm overflow-y-auto max-h-80"></div>
    </div>
  </div>

  <!-- ── REPORT PHASE ────────────────────────────────────────── -->
  <div id="ph-report" class="hidden">
    <!-- Toolbar -->
    <div class="no-print flex flex-wrap items-center gap-2 mb-4 surface rounded-xl p-3 shadow-sm sticky top-2 z-40">
      <button onclick="window.print()" class="btn btn-gray">🖨️ In/PDF</button>
      <button id="btn-slides" onclick="doSlides()" class="btn btn-gray">📑 Slides PPTX</button>
      <button id="btn-docx" onclick="doDocx()" class="btn btn-gray">📄 Word</button>
      <button onclick="openReports()" class="btn btn-gray">📂 Báo cáo đã lưu</button>
      <div class="flex-1"></div>
      <button onclick="chgFont(-1)" class="btn btn-gray font-bold">A−</button>
      <button onclick="chgFont(1)"  class="btn btn-gray font-bold">A+</button>
      <button onclick="toggleDark()" id="dark-btn2" class="btn btn-gray">☀️</button>
      <button onclick="newReport()" class="btn btn-green">+ Mới</button>
    </div>

    <!-- Report card -->
    <div class="surface rounded-xl shadow-sm overflow-hidden">
      <!-- Header band -->
      <div class="bg-green-700 text-white px-6 py-5">
        <p class="text-green-300 text-xs uppercase tracking-widest font-medium">Báo cáo phân tích thuế</p>
        <h1 id="rpt-title" class="text-xl font-bold mt-1">—</h1>
        <p id="rpt-meta" class="text-green-300 text-xs mt-1"></p>
      </div>
      <!-- TOC -->
      <div id="toc-wrap" class="border-b px-6 py-3 bg-gray-50 hidden" style="background:var(--bg)">
        <p class="text-xs font-semibold uppercase tracking-wide mb-2" style="color:var(--muted)">Mục lục</p>
        <div id="toc-list" class="space-y-0.5"></div>
      </div>
      <!-- Content -->
      <div id="report-content" class="px-8 py-6" style="font-size:16px"></div>
      <!-- Sources -->
      <div id="src-wrap" class="border-t px-6 py-4 hidden" style="background:var(--bg)">
        <button onclick="toggleSrc()" class="text-sm font-semibold flex items-center gap-2" style="color:var(--muted)">
          <span id="src-arrow">▶</span> Nguồn tham khảo (<span id="src-count">0</span>)
        </button>
        <div id="src-list" class="hidden mt-2 space-y-1"></div>
      </div>
    </div>
  </div>
</div><!-- /app -->

<!-- ── Reports Modal ─────────────────────────────────────────── -->
<div id="modal-reports" class="fixed inset-0 bg-black/50 z-50 hidden flex items-center justify-center">
  <div class="bg-white rounded-2xl shadow-2xl w-full max-w-xl mx-4 max-h-[85vh] flex flex-col"
       style="background:var(--surface);color:var(--text)">
    <!-- Header -->
    <div class="p-5 border-b flex items-center justify-between" style="border-color:var(--border)">
      <div>
        <h2 class="font-bold text-lg">📂 Báo cáo đã lưu</h2>
        <p class="text-xs mt-0.5" style="color:var(--muted)">Mới nhất → cũ nhất</p>
      </div>
      <button onclick="closeModal('modal-reports')" class="text-gray-400 hover:text-gray-600 text-xl leading-none">✕</button>
    </div>
    <!-- Search box -->
    <div class="px-4 pt-3 pb-2">
      <input type="text" id="reports-search"
        placeholder="🔍 Tìm kiếm theo tên báo cáo..."
        oninput="filterReports(this.value)"
        class="w-full px-3 py-2 text-sm rounded-lg border focus:outline-none focus:ring-2"
        style="border-color:var(--border);background:var(--bg);color:var(--text);--tw-ring-color:var(--accent)">
    </div>
    <!-- List -->
    <div id="reports-list" class="flex-1 overflow-y-auto px-4 pb-4 space-y-2"></div>
  </div>
</div>


<script>
// ── State ─────────────────────────────────────────────────────
let AUTH = '';
let mode = 'sector';
let sections = [];
let reportHtml = '';
let citations = [];
let currentFile = null;
let fontSize = 16;
let dark = false;
let abortCtrl = null;

// ── Auth ──────────────────────────────────────────────────────
async function doLogin() {
  const u = document.getElementById('li-user').value.trim();
  const p = document.getElementById('li-pass').value.trim();
  AUTH = 'Basic ' + btoa(u + ':' + p);
  try {
    const r = await fetch('/health', {headers: {Authorization: AUTH}});
    if (r.ok) {
      document.getElementById('login-modal').classList.add('hidden');
      init();
      // Load reports count for badge
      fetch('/reports', {headers: {Authorization: AUTH}})
        .then(r => r.json())
        .then(data => {
          const badge = document.getElementById('reports-badge');
          if (badge) badge.textContent = data.length ? `(${data.length})` : '';
        }).catch(() => {});
    } else {
      document.getElementById('li-err').classList.remove('hidden');
      AUTH = '';
    }
  } catch(e) {
    document.getElementById('li-err').classList.remove('hidden');
    AUTH = '';
  }
}

// ── Init ──────────────────────────────────────────────────────
async function init() {
  await loadSections();
}

// ── Mode ──────────────────────────────────────────────────────
function setMode(m) {
  mode = m;
  const on  = 'px-5 py-2 rounded-full text-sm font-medium transition bg-green-600 text-white';
  const off = 'px-5 py-2 rounded-full text-sm font-medium transition bg-gray-100 text-gray-600 hover:bg-gray-200';
  document.getElementById('tab-sector').className  = m === 'sector'  ? on : off;
  document.getElementById('tab-company').className = m === 'company' ? on : off;
  document.getElementById('subj-label').textContent = m === 'sector'
    ? 'Tên ngành / lĩnh vực cần phân tích:' : 'Tên công ty cần phân tích:';
  document.getElementById('subj-input').placeholder = m === 'sector'
    ? 'VD: Bất động sản, Fintech, Sản xuất thép...' : 'VD: Vinamilk, FPT, Masan Group...';
  loadSections();
}

// ── Sections ──────────────────────────────────────────────────
async function loadSections() {
  const r = await fetch('/default-sections?mode=' + mode, {headers: {Authorization: AUTH}});
  sections = await r.json();
  renderSections();
}

function resetSections() { loadSections(); }

function renderSections() {
  const el = document.getElementById('sections-list');
  el.innerHTML = '';
  sections.forEach((s, i) => el.appendChild(makeCard(s, i)));
}

function esc(s) {
  return String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
}

function makeCard(sec, i) {
  const div = document.createElement('div');
  div.className = 'sec-card' + (sec.enabled ? '' : ' off');
  div.dataset.id = sec.id;

  const chips = sec.sub.map((s, j) => `<span class="sub-chip">${esc(s)}<button onclick="rmSub('${sec.id}',${j})">✕</button></span>`).join('');

  div.innerHTML = `
    <div class="flex items-center gap-2 mb-2">
      <input type="checkbox" ${sec.enabled ? 'checked' : ''} class="accent-green-600 w-4 h-4 cursor-pointer"
        onchange="togSec('${sec.id}',this.checked)">
      <input type="text" value="${esc(sec.title)}"
        class="flex-1 font-medium text-sm border-0 focus:outline-none focus:ring-1 focus:ring-green-300 rounded px-1"
        style="background:transparent"
        onblur="updTitle('${sec.id}',this.value)">
      <button onclick="suggestSubs('${sec.id}')" class="text-xs px-2 py-0.5 rounded border border-gray-200 hover:bg-gray-50" style="color:var(--muted)" title="AI gợi ý chủ đề con">✨</button>
      <button onclick="rmSec('${sec.id}')" class="text-gray-300 hover:text-red-400 text-sm px-1">✕</button>
    </div>
    <div class="flex flex-wrap gap-1 ml-6">
      ${chips}
      <button onclick="addSub('${sec.id}')"
        class="text-xs px-2 py-0.5 rounded-full border border-dashed border-gray-300 hover:border-green-400 hover:text-green-600 transition"
        style="color:var(--muted)">+ thêm</button>
    </div>`;
  return div;
}

function togSec(id, enabled) {
  const s = sections.find(x => x.id === id);
  if (s) {
    s.enabled = enabled;
    const card = document.querySelector(`[data-id="${id}"]`);
    if (card) card.className = 'sec-card' + (enabled ? '' : ' off');
  }
}

function updTitle(id, t) {
  const s = sections.find(x => x.id === id);
  if (s) s.title = t;
}

function rmSec(id) {
  sections = sections.filter(s => s.id !== id);
  renderSections();
}

function addSection() {
  const id = 's' + Date.now();
  sections.push({id, title: 'Phần mới', sub: [], enabled: true});
  renderSections();
  setTimeout(() => {
    const inp = document.querySelector(`[data-id="${id}"] input[type="text"]`);
    if (inp) { inp.focus(); inp.select(); }
  }, 50);
}

function rmSub(secId, idx) {
  const s = sections.find(x => x.id === secId);
  if (s) { s.sub.splice(idx, 1); renderSections(); }
}

function addSub(secId) {
  const name = prompt('Nhập tên chủ đề con:');
  if (name && name.trim()) {
    const s = sections.find(x => x.id === secId);
    if (s) { s.sub.push(name.trim()); renderSections(); }
  }
}

async function suggestSubs(secId) {
  const s = sections.find(x => x.id === secId);
  if (!s) return;
  const subject = document.getElementById('subj-input').value.trim() || '(chưa nhập)';
  const r = await fetch('/suggest-subsections', {
    method: 'POST',
    headers: {Authorization: AUTH, 'Content-Type': 'application/json'},
    body: JSON.stringify({title: s.title, subject})
  });
  const {suggestions} = await r.json();
  if (suggestions && suggestions.length) {
    suggestions.forEach(sg => { if (!s.sub.includes(sg)) s.sub.push(sg); });
    renderSections();
  }
}

function buildRefsHtml() {
  const srcList = document.getElementById('src-list');
  if (!srcList) return '';
  const items = srcList.querySelectorAll('a, li');
  if (!items.length) return '';
  let html = '<h2>Nguồn tham khảo</h2><ol>';
  items.forEach(el => {
    const a = el.tagName === 'A' ? el : el.querySelector('a');
    if (a) {
      html += `<li><a href="${a.href}" target="_blank">${a.textContent.trim()}</a></li>`;
    } else {
      const t = el.textContent.trim();
      if (t) html += `<li>${t}</li>`;
    }
  });
  html += '</ol>';
  return html;
}

// ── Research ──────────────────────────────────────────────────

async function startResearch() {
  const subject = document.getElementById('subj-input').value.trim();
  if (!subject) { alert('Vui lòng nhập tên ngành hoặc công ty'); return; }
  const enabled = sections.filter(s => s.enabled);
  if (!enabled.length) { alert('Vui lòng bật ít nhất một phần'); return; }

  const sonar = document.querySelector('input[name="sonar"]:checked').value;
  showPhase('progress');
  document.getElementById('prog-subject').textContent = subject;
  reportHtml = '';
  citations = [];
  currentFile = null;
  document.getElementById('step-list').innerHTML = '';
  document.getElementById('partial-body').innerHTML = '';
  document.getElementById('partial-wrap').classList.add('hidden');

  abortCtrl = new AbortController();

  try {
    const res = await fetch('/stream', {
      method: 'POST',
      headers: {Authorization: AUTH, 'Content-Type': 'application/json'},
      body: JSON.stringify({subject, mode, sections, sonar_model: sonar}),
      signal: abortCtrl.signal,
    });

    if (!res.ok) { showErr('Lỗi server: ' + res.status); return; }

    const reader = res.body.getReader();
    const dec = new TextDecoder();
    let buf = '';

    while (true) {
      const {done, value} = await reader.read();
      if (done) break;
      buf += dec.decode(value, {stream: true});
      const lines = buf.split('\n');
      buf = lines.pop();
      for (const line of lines) {
        if (line.startsWith('data: ')) {
          try { handleEvt(JSON.parse(line.slice(6)), subject); } catch(e) {}
        }
      }
    }
  } catch(e) {
    if (e.name !== 'AbortError') showErr('Lỗi kết nối: ' + e.message);
  }
}

function cancelResearch() {
  if (abortCtrl) { abortCtrl.abort(); abortCtrl = null; }
  showPhase('setup');
}

function handleEvt(evt, subject) {
  switch (evt.type) {
    case 'status':
      document.getElementById('prog-status').textContent = evt.message;
      break;
    case 'progress': {
      const pct = Math.round(evt.step / evt.total * 100);
      document.getElementById('prog-bar').style.width = pct + '%';
      document.getElementById('prog-pct').textContent = pct + '%';
      document.getElementById('prog-label').textContent = evt.label || '';
      const item = document.createElement('div');
      item.className = 'flex items-center gap-2 text-sm surface rounded-lg px-3 py-2 fade-in';
      item.innerHTML = '<span class="text-green-500">✓</span> ' + esc(evt.label || '');
      document.getElementById('step-list').appendChild(item);
      item.scrollIntoView({behavior: 'smooth', block: 'nearest'});
      break;
    }
    case 'ai_start':
      document.getElementById('prog-status').textContent = 'AI đang viết báo cáo...';
      document.getElementById('partial-wrap').classList.remove('hidden');
      break;
    case 'chunk':
      reportHtml += evt.text;
      document.getElementById('partial-body').innerHTML += evt.text;
      break;
    case 'citations':
      citations = citations.concat(evt.urls);
      break;
    case 'done':
      currentFile = evt.filename;
      finishReport(subject);
      break;
    case 'error':
      showErr(evt.message);
      break;
  }
}

function finishReport(subject) {
  showPhase('report');
  document.getElementById('rpt-title').textContent = 'Phân Tích Thuế — ' + subject;
  document.getElementById('rpt-meta').textContent =
    'Ngày tạo: ' + new Date().toLocaleDateString('vi-VN') +
    (currentFile ? ' · Đã lưu: ' + currentFile : '');

  const content = document.getElementById('report-content');
  content.innerHTML = reportHtml;
  content.style.fontSize = fontSize + 'px';

  // Add AI caveat at bottom
  const caveat = document.createElement('div');
  caveat.style.cssText = 'margin-top:2rem;padding:1rem;background:var(--bg);border-top:2px solid var(--border);border-radius:.5rem;font-size:.8rem;color:var(--muted);line-height:1.6';
  caveat.innerHTML = '<strong>&#9888;&#65039; Lưu ý:</strong> Báo cáo tạo bởi AI (Tax Sector Research). Mang tính tham khảo, không thay thế tư vấn thuế chuyên nghiệp. Kiểm chứng hiệu lực văn bản tại <a href="https://thuvienphapluat.vn" target="_blank" style="color:var(--brand)">thuvienphapluat.vn</a>.';
  content.appendChild(caveat);

  buildTOC();
  buildSources();
  setupScroll();
  window.scrollTo({top: 0, behavior: 'smooth'});
}

// ── TOC ───────────────────────────────────────────────────────
function buildTOC() {
  const h2s = document.getElementById('report-content').querySelectorAll('h2');
  if (h2s.length < 2) return;
  const list = document.getElementById('toc-list');
  list.innerHTML = '';
  h2s.forEach((h, i) => {
    h.id = 'sec-' + i;
    const a = document.createElement('a');
    a.href = '#sec-' + i;
    a.className = 'block text-sm py-0.5 hover:underline';
    a.style.color = 'var(--brand)';
    a.textContent = h.textContent;
    a.onclick = e => { e.preventDefault(); h.scrollIntoView({behavior: 'smooth'}); };
    list.appendChild(a);
  });
  document.getElementById('toc-wrap').classList.remove('hidden');
}

// ── Sources ───────────────────────────────────────────────────
function buildSources() {
  const unique = [...new Set(citations)];
  if (!unique.length) return;
  const list = document.getElementById('src-list');
  list.innerHTML = '';
  unique.forEach((url, i) => {
    const d = document.createElement('div');
    d.className = 'text-xs';
    d.style.color = 'var(--muted)';
    d.innerHTML = `[${i+1}] <a href="${esc(url)}" target="_blank" class="hover:underline" style="color:var(--brand)">${esc(url)}</a>`;
    list.appendChild(d);
  });
  document.getElementById('src-count').textContent = unique.length;
  document.getElementById('src-wrap').classList.remove('hidden');
}

function toggleSrc() {
  const list = document.getElementById('src-list');
  const hidden = list.classList.toggle('hidden');
  document.getElementById('src-arrow').textContent = hidden ? '▶' : '▼';
}

// ── Scroll progress ───────────────────────────────────────────
function setupScroll() {
  window.addEventListener('scroll', () => {
    const d = document.documentElement;
    const pct = d.scrollTop / (d.scrollHeight - d.clientHeight);
    document.getElementById('reading-bar').style.width = (pct * 100) + '%';
  });
}

// ── Toolbar ───────────────────────────────────────────────────
function chgFont(delta) {
  fontSize = Math.max(12, Math.min(22, fontSize + delta));
  document.getElementById('report-content').style.fontSize = fontSize + 'px';
}

function toggleDark() {
  dark = !dark;
  document.body.classList.toggle('dark', dark);
  const icon = dark ? '🌙' : '☀️';
  document.getElementById('dark-btn').textContent = icon;
  const b2 = document.getElementById('dark-btn2');
  if (b2) b2.textContent = icon;
}

async function doSlides() {
  if (!reportHtml) return;
  const subject = document.getElementById('rpt-title').textContent.replace('Phân Tích Thuế — ', '');
  const btn = document.getElementById('btn-slides');
  btn.textContent = '⏳ Đang tạo PPTX...';
  btn.disabled = true;
  try {
    const r = await fetch('/slides', {
      method: 'POST',
      headers: {Authorization: AUTH, 'Content-Type': 'application/json'},
      body: JSON.stringify({html: reportHtml, subject}),
    });
    if (r.ok) {
      const blob = await r.blob();
      const url = URL.createObjectURL(blob);
      const a = document.createElement('a');
      a.href = url;
      a.download = (currentFile || subject).replace('.html', '') + '.pptx';
      a.click();
      URL.revokeObjectURL(url);
    } else {
      alert('Không thể tạo PPTX');
    }
  } catch(e) { alert('Lỗi: ' + e.message); }
  finally { btn.textContent = '📑 Slides PPTX'; btn.disabled = false; }
}

async function doDocx() {
  const btn = document.getElementById('btn-docx');
  if (!reportHtml) { alert('Chưa có báo cáo để xuất'); return; }
  const originalText = btn ? btn.textContent : '';
  if (btn) { btn.textContent = '⏳ Đang xuất...'; btn.disabled = true; }
  try {
    const controller = new AbortController();
    const timeoutId = setTimeout(() => controller.abort(), 60000);
    const r = await fetch('/docx', {
      method: 'POST',
      headers: {Authorization: AUTH, 'Content-Type': 'application/json'},
      body: JSON.stringify({
        html: reportHtml,
        subject: document.getElementById('rpt-title')?.textContent || 'Báo cáo'
      }),
      signal: controller.signal,
    });
    clearTimeout(timeoutId);
    if (!r.ok) {
      const err = await r.json().catch(() => ({detail: 'Lỗi không xác định'}));
      throw new Error(err.detail || `HTTP ${r.status}`);
    }
    const blob = await r.blob();
    const url = URL.createObjectURL(blob);
    const a = document.createElement('a');
    a.href = url;
    a.download = `PhanTichThue_${(document.getElementById('rpt-title')?.textContent || 'BaoCao').replace(/\s+/g,'_').slice(0,50)}.docx`;
    document.body.appendChild(a);
    a.click();
    document.body.removeChild(a);
    URL.revokeObjectURL(url);
  } catch(e) {
    if (e.name === 'AbortError') {
      alert('Xuất DOCX quá thời gian (60s). Báo cáo có thể quá dài — thử lại hoặc rút ngắn nội dung.');
    } else {
      alert('Xuất DOCX thất bại: ' + e.message);
    }
  } finally {
    if (btn) { btn.textContent = originalText; btn.disabled = false; }
  }
}

// ── Reports modal ─────────────────────────────────────────────

async function openReports() {
  openModal('modal-reports');
  const list = document.getElementById('reports-list');
  list.innerHTML = '<p class="text-sm text-center py-4" style="color:var(--muted)">Đang tải...</p>';
  const searchEl = document.getElementById('reports-search');
  if (searchEl) searchEl.value = '';
  try {
    const r = await fetch('/reports', {headers: {Authorization: AUTH}});
    if (!r.ok) throw new Error('Lỗi tải danh sách');
    const data = await r.json();
    allReports = data.sort((a, b) => (b.mtime || 0) - (a.mtime || 0));
    const badge = document.getElementById('reports-badge');
    if (badge) badge.textContent = allReports.length ? `(${allReports.length})` : '';
    renderReports(allReports, false);
  } catch(e) {
    list.innerHTML = `<p class="text-sm text-center py-4 text-red-500">Lỗi: ${e.message}</p>`;
  }
}

function renderReports(list, showAll) {
  const container = document.getElementById('reports-list');
  if (!list.length) {
    container.innerHTML = '<p class="text-sm py-4 text-center" style="color:var(--muted)">Chưa có báo cáo nào</p>';
    return;
  }
  const LIMIT = 10;
  const displayed = showAll ? list : list.slice(0, LIMIT);
  const remaining = list.length - LIMIT;
  let html = displayed.map(f => {
    const name = f.name || f;
    const subject = f.subject || esc(name.replace('.html',''));
    const date = f.mtime ? new Date(f.mtime * 1000).toLocaleDateString('vi-VN', {
      day:'2-digit', month:'2-digit', year:'numeric', hour:'2-digit', minute:'2-digit'
    }) : (f.date || '');
    const size = f.size ? (f.size > 1024*1024 ? (f.size/1024/1024).toFixed(1)+'MB' : Math.round(f.size/1024)+'KB') : '';
    return `
      <div class="flex items-center gap-2 p-3 rounded-lg border fade-in group" style="border-color:var(--border)">
        <span class="text-base shrink-0">📄</span>
        <div class="flex-1 min-w-0 cursor-pointer" onclick="loadSavedReport('${esc(name)}')">
          <p class="font-medium text-sm truncate">${subject}</p>
          ${date ? `<p class="text-xs" style="color:var(--muted)">${date}${size ? ' · '+size : ''}</p>` : ''}
        </div>
        <button onclick="loadSavedReport('${esc(name)}')" class="btn btn-green text-xs shrink-0">📂 Mở</button>
        <a href="/report/${encodeURIComponent(name)}" target="_blank" class="btn btn-gray text-xs shrink-0">↗</a>
        <button onclick="delReport('${esc(name)}',this)" class="btn text-xs shrink-0 text-red-400 hover:bg-red-50">🗑</button>
      </div>`;
  }).join('');
  if (!showAll && remaining > 0) {
    html += `<button onclick="renderReports(allReports, true)"
      class="w-full mt-2 py-2 text-sm rounded-lg border" style="border-color:var(--border);color:var(--accent)">
      ↓ Xem thêm ${remaining} báo cáo</button>`;
  }
  container.innerHTML = html;
}

function filterReports(query) {
  const q = query.toLowerCase().trim();
  if (!q) { renderReports(allReports, false); return; }
  const filtered = allReports.filter(f => {
    const name = (f.name || f).toLowerCase();
    const subject = (f.subject || '').toLowerCase();
    return name.includes(q) || subject.includes(q);
  });
  renderReports(filtered, true);
}

async function loadSavedReport(name) {
  closeModal('modal-reports');
  try {
    const r = await fetch('/report/' + encodeURIComponent(name), {headers: {Authorization: AUTH}});
    const html = await r.text();
    const parser = new DOMParser();
    const doc = parser.parseFromString(html, 'text/html');
    const body = doc.getElementById('report-body');
    reportHtml = body ? body.innerHTML : doc.body.innerHTML;
    currentFile = name;
    citations = [];

    // Try to extract title
    const h1 = doc.querySelector('h1');
    const subject = h1 ? h1.textContent.replace('Phan Tich Thue — ', '').replace('Phân Tích Thuế — ', '') : name;

    showPhase('report');
    document.getElementById('rpt-title').textContent = 'Phân Tích Thuế — ' + subject;
    document.getElementById('rpt-meta').textContent = 'Đã lưu: ' + name;
    document.getElementById('report-content').innerHTML = reportHtml;
    document.getElementById('report-content').style.fontSize = fontSize + 'px';
    document.getElementById('toc-wrap').classList.add('hidden');
    document.getElementById('src-wrap').classList.add('hidden');
    buildTOC();
    setupScroll();
    window.scrollTo({top: 0, behavior: 'smooth'});
  } catch(e) { alert('Không thể tải báo cáo: ' + e.message); }
}

async function delReport(name, btn) {
  if (!confirm('Xóa báo cáo "' + name + '"?')) return;
  try {
    const r = await fetch('/report/' + encodeURIComponent(name), {
      method: 'DELETE', headers: {Authorization: AUTH}
    });
    if (r.ok) btn.closest('.flex').remove();
    else alert('Không thể xóa');
  } catch(e) { alert('Lỗi: ' + e.message); }
}

function newReport() {
  if (currentFile || reportHtml) {
    if (!confirm('Bắt đầu báo cáo mới? Báo cáo hiện tại vẫn được lưu.')) return;
  }
  showPhase('setup');
  reportHtml = ''; currentFile = null; citations = [];
  document.getElementById('step-list').innerHTML = '';
  document.getElementById('partial-body').innerHTML = '';
  document.getElementById('report-content').innerHTML = '';
  document.getElementById('toc-wrap').classList.add('hidden');
  document.getElementById('src-wrap').classList.add('hidden');
  window.scrollTo({top: 0, behavior: 'smooth'});
}

// ── Helpers ───────────────────────────────────────────────────
function showPhase(p) {
  ['setup','progress','report'].forEach(x =>
    document.getElementById('ph-' + x).classList.toggle('hidden', x !== p));
}

function showErr(msg) {
  document.getElementById('prog-status').textContent = '❌ ' + msg;
  document.getElementById('prog-bar').style.background = '#ef4444';
}

function openModal(id)  { document.getElementById(id).classList.remove('hidden'); }
function closeModal(id) { document.getElementById(id).classList.add('hidden'); }

// Close modal on backdrop click
document.querySelectorAll('[id^="modal-"]').forEach(m =>
  m.addEventListener('click', e => { if (e.target === m) closeModal(m.id); }));

document.addEventListener('keydown', e => {
  if (e.key === 'Escape') closeModal('modal-reports');
});
</script>
</body>
</html>"""

# ── Favicon ───────────────────────────────────────────────────────────────────
@app.get("/favicon.svg", include_in_schema=False)
async def favicon():
    import pathlib
    svg_path = pathlib.Path(__file__).parent / "favicon.svg"
    if svg_path.exists():
        return FileResponse(svg_path, media_type="image/svg+xml")
    return Response(status_code=404)

# ── Serve frontend ────────────────────────────────────────────────────────────
@app.get("/", response_class=HTMLResponse)
def root():
    return HTMLResponse(HTML_PAGE)
